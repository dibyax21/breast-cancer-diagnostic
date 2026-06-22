import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Set widescreen 16:9 format
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Styling Colors
BG_DARK = RGBColor(5, 12, 22)
CARD_DARK = RGBColor(15, 28, 46)
GUIDE_DARK = RGBColor(10, 18, 30)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(0, 190, 214)
PURPLE = RGBColor(186, 85, 211)
GREY_TEXT = RGBColor(160, 174, 192)

# Helper to configure slide background
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

# Helper to add standard slide footer
def add_footer(slide, current, total=10, center_text="MODEL VERIFICATION DECK"):
    # Left: ONCOLOGY.AI
    tx_left = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(3.0), Inches(0.4))
    tf_l = tx_left.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.text = "\u2738 ONCOLOGY.AI"
    p_l.font.bold = True
    p_l.font.size = Pt(10)
    p_l.font.color.rgb = TEAL
    p_l.font.name = "Arial"
    
    # Center: Slide Descriptor
    tx_center = slide.shapes.add_textbox(Inches(4.0), Inches(6.9), Inches(5.33), Inches(0.4))
    tf_c = tx_center.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    p_c.text = center_text.upper()
    p_c.font.size = Pt(9)
    p_c.font.color.rgb = GREY_TEXT
    p_c.font.name = "Arial"
    
    # Right: SLIDE X OF Y
    tx_right = slide.shapes.add_textbox(Inches(9.8), Inches(6.9), Inches(2.7), Inches(0.4))
    tf_r = tx_right.text_frame
    p_r = tf_r.paragraphs[0]
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.text = f"SLIDE {current} OF {total}"
    p_r.font.size = Pt(10)
    p_r.font.color.rgb = WHITE
    p_r.font.name = "Arial"

# Helper to add slide title (w/ dual color support)
def add_slide_title(slide, part_white, part_teal):
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    
    # White part
    run1 = p.add_run()
    run1.text = part_white + " "
    run1.font.bold = True
    run1.font.size = Pt(28)
    run1.font.color.rgb = WHITE
    run1.font.name = "Arial"
    
    # Teal part
    run2 = p.add_run()
    run2.text = part_teal
    run2.font.bold = True
    run2.font.size = Pt(28)
    run2.font.color.rgb = TEAL
    run2.font.name = "Arial"

# Helper to add Presenter Guide card on the right
def add_presenter_guide(slide, phase_text, title_text, bullets, quote_text):
    # Base container
    card = slide.shapes.add_shape(1, Inches(9.8), Inches(1.2), Inches(2.7), Inches(5.2))
    card.fill.solid()
    card.fill.fore_color.rgb = GUIDE_DARK
    card.line.color.rgb = CARD_DARK
    card.line.width = Pt(1.5)
    
    # Header Pill Tag: PRESENTER GUIDE
    tag = slide.shapes.add_shape(1, Inches(10.0), Inches(0.95), Inches(1.2), Inches(0.28))
    tag.fill.solid()
    tag.fill.fore_color.rgb = TEAL
    tag.line.fill.background()
    p_tag = tag.text_frame.paragraphs[0]
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.text = "PRESENTER GUIDE"
    p_tag.font.bold = True
    p_tag.font.size = Pt(8)
    p_tag.font.color.rgb = BG_DARK
    p_tag.font.name = "Arial"
    
    # Text Frame for content
    tx_box = slide.shapes.add_textbox(Inches(9.95), Inches(1.35), Inches(2.4), Inches(2.6))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    # Phase text
    p_phase = tf.paragraphs[0]
    p_phase.text = phase_text.upper()
    p_phase.font.bold = True
    p_phase.font.size = Pt(9)
    p_phase.font.color.rgb = PURPLE
    p_phase.font.name = "Arial"
    p_phase.space_after = Pt(4)
    
    # Title text
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.bold = True
    p_title.font.size = Pt(14)
    p_title.font.color.rgb = WHITE
    p_title.font.name = "Arial"
    p_title.space_after = Pt(12)
    
    # Bullet points
    for bullet in bullets:
        p_b = tf.add_paragraph()
        p_b.text = bullet
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEAL
        p_b.font.name = "Arial"
        p_b.space_after = Pt(8)
        
    # Quote Card at the bottom
    # Draw left border line
    line = slide.shapes.add_shape(1, Inches(10.05), Inches(4.2), Inches(0.04), Inches(1.8))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    
    # Quote text box
    q_box = slide.shapes.add_textbox(Inches(10.15), Inches(4.15), Inches(2.2), Inches(1.8))
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    p_q = q_tf.paragraphs[0]
    p_q.text = f'"{quote_text}"'
    p_q.font.italic = True
    p_q.font.size = Pt(11)
    p_q.font.color.rgb = GREY_TEXT
    p_q.font.name = "Arial"

# Helper to add standard text card
def add_content_card(slide, left, top, width, height, title, items):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_DARK
    card.line.fill.background()
    
    tx_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = GREY_TEXT
        p.font.name = "Arial"
        p.space_after = Pt(10)
        
        # Format bold titles within cards
        if ":" in item:
            parts = item.split(":", 1)
            p.text = ""
            r_bold = p.add_run()
            r_bold.text = parts[0] + ":"
            r_bold.font.bold = True
            r_bold.font.color.rgb = WHITE
            
            r_normal = p.add_run()
            r_normal.text = parts[1]
            r_normal.font.color.rgb = GREY_TEXT
    return card

# Helper to add an image card (like the screenshots)
def add_visual_card(slide, image_path, left, top, width, height, label="MODEL ARTIFACT SCAN"):
    # Border/container shape
    container = slide.shapes.add_shape(1, left, top, width, height)
    container.fill.solid()
    container.fill.fore_color.rgb = CARD_DARK
    container.line.fill.background()
    
    # Add image inside container (leave a bit of padding for border)
    pad = Inches(0.1)
    slide.shapes.add_picture(image_path, left + pad, top + pad, width - pad*2, height - pad*2)
    
    # Add small bottom label tag
    tag = slide.shapes.add_shape(1, left + pad, top + height - Inches(0.35), Inches(1.3), Inches(0.22))
    tag.fill.solid()
    tag.fill.fore_color.rgb = TEAL
    tag.line.fill.background()
    
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.bold = True
    p.font.size = Pt(7)
    p.font.color.rgb = BG_DARK
    p.font.name = "Arial"


# ==========================================
# SLIDE 1: Title Slide (Dark Theme)
# ==========================================
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)
set_dark_background(slide1)

# Pill Tag
tag1 = slide1.shapes.add_shape(1, Inches(2.0), Inches(2.0), Inches(2.6), Inches(0.35))
tag1.fill.solid()
tag1.fill.fore_color.rgb = GUIDE_DARK
tag1.line.color.rgb = TEAL
tag1.line.width = Pt(1.5)
p_tag1 = tag1.text_frame.paragraphs[0]
p_tag1.text = "\u2699 CLINICAL ML DEPLOYMENT PROGRAM"
p_tag1.font.bold = True
p_tag1.font.size = Pt(9)
p_tag1.font.color.rgb = TEAL
p_tag1.font.name = "Arial"

# Title Box
title_box = slide1.shapes.add_textbox(Inches(2.0), Inches(2.6), Inches(9.33), Inches(3.0))
tf1 = title_box.text_frame
tf1.word_wrap = True

# Main Title
p = tf1.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run1 = p.add_run()
run1.text = "AI-Driven Breast Cancer\n"
run1.font.bold = True
run1.font.size = Pt(44)
run1.font.color.rgb = WHITE
run1.font.name = "Arial"

run2 = p.add_run()
run2.text = "Diagnostic Portal"
run2.font.bold = True
run2.font.size = Pt(44)
run2.font.color.rgb = TEAL
run2.font.name = "Arial"

# Subtitle
p2 = tf1.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "Parallel Machine Learning & Deep Learning Decision Support Clinical Classification System. Featuring empirical model outputs validation metrics."
p2.font.size = Pt(16)
p2.font.color.rgb = GREY_TEXT
p2.font.name = "Arial"
p2.space_before = Pt(18)

# Reference note
p3 = tf1.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run_ref = p3.add_run()
run_ref.text = "Baseline Repository References: "
run_ref.font.color.rgb = GREY_TEXT
run_ref.font.size = Pt(11)
run_ref.font.name = "Arial"

run_ref2 = p3.add_run()
run_ref2.text = "breast_cancer_diagnostic_presentation_v2.pptx"
run_ref2.font.bold = True
run_ref2.font.color.rgb = WHITE
run_ref2.font.size = Pt(11)
run_ref2.font.name = "Arial"
p3.space_before = Pt(30)

add_footer(slide1, 1, center_text="MODEL VERIFICATION DECK")


# ==========================================
# SLIDE 2: Oncology Context & Challenge
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
set_dark_background(slide2)
add_slide_title(slide2, "Oncology Context &", "Diagnostic Challenge")

# Left content card
add_content_card(slide2, Inches(0.8), Inches(2.0), Inches(4.0), Inches(4.2), "", [
    "\u26A0 Critical Mortality: Breast cancer remains a primary driver of global oncology mortality; rapid, highly accurate diagnosis saves lives.",
    "\u2695 Cytological FNA Analysis: Fine needle aspiration (FNA) biopsied cells are evaluated based on complex structural characteristics.",
    "\u2714 Objective & Redundancy: Classify cell samples safely into Benign or Malignant classes using redundant parallel baselines."
])

# Middle image card
add_visual_card(slide2, "fna_scan.jpg", Inches(5.0), Inches(2.0), Inches(4.3), Inches(4.2))

# Presenter guide
add_presenter_guide(slide2, "Phase 1: Intake", "Intake & Redundancy", [
    "\u2192 FNA biopsy limits error.",
    "\u2192 Automating histology splits clinical risk."
], "We automate cytological categorization because fine-needle analysis yields complex, noisy cell structures.")

add_footer(slide2, 2, center_text="CLINICAL LANDSCAPE")


# ==========================================
# SLIDE 3: System Pipeline & Data Flow
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
set_dark_background(slide3)
add_slide_title(slide3, "System Pipeline &", "Data Flow Architecture")

# Card 1: Gold Ingestion
add_content_card(slide3, Inches(0.8), Inches(2.5), Inches(2.7), Inches(3.4), "", [
    "\U0001F4C4 01 / Gold Ingestion",
    "Ingests Wisconsin clinical repository dataset vectors containing cell structural attributes."
])

# Card 2: Standardization
card2 = add_content_card(slide3, Inches(3.7), Inches(2.5), Inches(2.7), Inches(3.4), "", [
    "\u2699 02 / Standardization",
    "Resolves missing elements marked '?' and fits features to normal distribution bounds."
])
# Add Formula Z = (x - mu) / sigma
# Container box
box2 = slide3.shapes.add_shape(1, Inches(4.1), Inches(4.5), Inches(1.9), Inches(0.9))
box2.fill.solid()
box2.fill.fore_color.rgb = GUIDE_DARK
box2.line.color.rgb = CARD_DARK
box2.line.width = Pt(1.5)

# Text Z =
z_box = slide3.shapes.add_textbox(Inches(4.2), Inches(4.7), Inches(0.5), Inches(0.4))
z_box.text_frame.paragraphs[0].text = "Z = "
z_box.text_frame.paragraphs[0].font.size = Pt(14)
z_box.text_frame.paragraphs[0].font.color.rgb = TEAL
z_box.text_frame.paragraphs[0].font.bold = True

# Numerator x - mu
num_box = slide3.shapes.add_textbox(Inches(4.8), Inches(4.45), Inches(1.0), Inches(0.3))
num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
num_box.text_frame.paragraphs[0].text = "x - \u03bc"
num_box.text_frame.paragraphs[0].font.size = Pt(12)
num_box.text_frame.paragraphs[0].font.color.rgb = WHITE

# Division Line
div_line = slide3.shapes.add_shape(1, Inches(4.8), Inches(4.88), Inches(0.9), Inches(0.02))
div_line.fill.solid()
div_line.fill.fore_color.rgb = WHITE
div_line.line.fill.background()

# Denominator sigma
den_box = slide3.shapes.add_textbox(Inches(4.8), Inches(4.95), Inches(1.0), Inches(0.3))
den_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
den_box.text_frame.paragraphs[0].text = "\u03c3"
den_box.text_frame.paragraphs[0].font.size = Pt(12)
den_box.text_frame.paragraphs[0].font.color.rgb = WHITE

# Card 3: Stratified Splitting
add_content_card(slide3, Inches(6.6), Inches(2.5), Inches(2.7), Inches(3.4), "", [
    "\u25d4 03 / Stratified Splitting",
    "Partitions features via an 80/20 train/test framework, safeguarding categorical balances."
])

# Presenter guide
add_presenter_guide(slide3, "Phase 2: Pipeline", "Data Preprocessing", [
    "\u2192 Clean data prevents skew.",
    "\u2192 StandardScaler prepares vector alignments."
], "The input variables are standardized to remove scales, preparing clinical inputs for dual classifiers.")

add_footer(slide3, 3, center_text="PREPROCESSING ARCHITECTURE")


# ==========================================
# SLIDE 4: Parallel Dual-Model Framework
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
set_dark_background(slide4)
add_slide_title(slide4, "Parallel Dual-Model", "Inference Pipeline")

# Card 1: Classifier Baseline
add_content_card(slide4, Inches(0.8), Inches(2.3), Inches(4.2), Inches(3.8), "", [
    "\u25d5 Classifier Baseline",
    "Random Forest Ensemble: Produces robust structural decisions to resolve high-variance outputs.",
    "Gini splits determine node branching thresholds."
])
# Card 1 Formula
box_g = slide4.shapes.add_shape(1, Inches(1.1), Inches(4.7), Inches(3.6), Inches(0.9))
box_g.fill.solid()
box_g.fill.fore_color.rgb = GUIDE_DARK
box_g.line.color.rgb = CARD_DARK
box_g.line.width = Pt(1.5)
g_tf = box_g.text_frame
p_g = g_tf.paragraphs[0]
p_g.alignment = PP_ALIGN.CENTER
p_g.text = "G = 1 - \u2211 P_i\u00b2"
p_g.font.size = Pt(16)
p_g.font.color.rgb = TEAL
p_g.font.bold = True
p_g.font.name = "Georgia"

# Card 2: Neural Engine Core
add_content_card(slide4, Inches(5.3), Inches(2.3), Inches(4.2), Inches(3.8), "", [
    "\u2699 Neural Engine Core",
    "Artificial Neural Net (ANN): Computes deep probability vectors mapping high-dimensional spaces.",
    "Includes local scikit-learn standard failover features."
])
# Card 2 Badge
badge = slide4.shapes.add_shape(1, Inches(5.6), Inches(4.7), Inches(3.6), Inches(0.6))
badge.fill.solid()
badge.fill.fore_color.rgb = GUIDE_DARK
badge.line.color.rgb = PURPLE
badge.line.width = Pt(1.5)
b_tf = badge.text_frame
p_b = b_tf.paragraphs[0]
p_b.alignment = PP_ALIGN.CENTER
p_b.text = "\u2738 BACKPROPAGATION ACTIVE"
p_b.font.bold = True
p_b.font.size = Pt(11)
p_b.font.color.rgb = PURPLE
p_b.font.name = "Arial"

# Presenter guide
add_presenter_guide(slide4, "Phase 3: Inference", "Dual Model Safety", [
    "\u2192 Parallel pathways reduce risk.",
    "\u2192 Contradictory classes trigger human path review."
], "We execute two diverse models side-by-side to guarantee classification redundancy and protect patient diagnostics.")

add_footer(slide4, 4, center_text="MODEL ARCHITECTURE")


# ==========================================
# SLIDE 5: Machine Learning Baseline (Random Forest)
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
set_dark_background(slide5)
add_slide_title(slide5, "Traditional ML Baseline:", "Random Forest")

# Left large accuracy metric
tx_acc = slide5.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.0), Inches(1.8))
tf_acc = tx_acc.text_frame
tf_acc.word_wrap = True
p_num = tf_acc.paragraphs[0]
p_num.alignment = PP_ALIGN.CENTER
p_num.text = "95.62%"
p_num.font.bold = True
p_num.font.size = Pt(84)
p_num.font.color.rgb = TEAL
p_num.font.name = "Arial Black"

p_label = tf_acc.add_paragraph()
p_label.alignment = PP_ALIGN.CENTER
p_label.text = "VALIDATION ACCURACY"
p_label.font.bold = True
p_label.font.size = Pt(14)
p_label.font.color.rgb = WHITE
p_label.font.name = "Arial"

# Right Card details
add_content_card(slide5, Inches(5.0), Inches(2.2), Inches(4.5), Inches(4.0), "", [
    "\u2192 Forest Composition: Employs 100 randomized decision estimators with entropy-driven entropy division metrics.",
    "\u2192 Anti-Overfitting Rules: Extreme resistance to outlier elements through features subset isolation.",
    "\u2192 High Traceability: Yields clear feature importances, enabling doctors to understand classification pathways."
])

# Presenter guide
add_presenter_guide(slide5, "Phase 4: Baseline", "Random Forest Bench", [
    "\u2192 95.62% test split result.",
    "\u2192 Establishes baseline clinical performance context."
], "Our Random Forest baseline provides a highly interpretable and strong benchmark at 95.62% out-of-sample accuracy.")

add_footer(slide5, 5, center_text="RANDOM FOREST BASE")


# ==========================================
# SLIDE 6: Performance Comparison (With Graphics)
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
set_dark_background(slide6)
add_slide_title(slide6, "Model Performance", "Comparison Metrics")

# Left Card
add_content_card(slide6, Inches(0.8), Inches(2.0), Inches(4.0), Inches(4.2), "", [
    "\u2699 Base Accuracy (RF): Achieved an initial baseline metric of 95.62%.",
    "\u26A1 Deep Learning ANN: Outperforms at 96.35% using dense layers and mathematical backpropagation.",
    "\u2714 Combined Trust: Executing models concurrently builds clinical confidence in positive tumor detections."
])

# Middle image card
add_visual_card(slide6, "comparison.png", Inches(5.0), Inches(2.0), Inches(4.3), Inches(4.2))

# Presenter guide
add_presenter_guide(slide6, "Phase 5: Evaluation", "Comparative Analysis", [
    "\u2192 Neural nets achieve top score.",
    "\u2192 Dual validation prevents failures."
], "Your run results show that while Random Forest is extremely strong, the Deep Learning topology captures complex, non-linear signals better.")

add_footer(slide6, 6, center_text="PERFORMANCE RUN")


# ==========================================
# SLIDE 7: Evaluation Toolkit (With Graphics)
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
set_dark_background(slide7)
add_slide_title(slide7, "ROC Curves &", "Confusion Matrix Verification")

# Left image card
add_visual_card(slide7, "evaluation_toolkit.png", Inches(0.8), Inches(2.0), Inches(4.3), Inches(4.2))

# Middle Card
add_content_card(slide7, Inches(5.4), Inches(2.0), Inches(4.1), Inches(4.2), "", [
    "\u2B50 AUC Ratio Analysis: Random Forest achieves 0.9926 while Neural Network maps 0.9859.",
    "\u2714 Zero False Negatives: The deep network mapped zero false negative samples on verification runs.",
    "\u2696 Balanced Thresholds: Keeps patient stress low by containing false positive instances."
])

# Presenter guide
add_presenter_guide(slide7, "Phase 6: Safety", "Patient Risk Mitigation", [
    "\u2192 High specificity levels.",
    "\u2192 Zero false negatives verified."
], "In clinical diagnostics, a false negative is deadly. Our ANN's zero false-negative output is a critical safeguard.")

add_footer(slide7, 7, center_text="ROC/CONFUSION VERIFICATION")


# ==========================================
# SLIDE 8: Explainable AI (With Graphics)
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
set_dark_background(slide8)
add_slide_title(slide8, "Explainable AI:", "Clinical Feature Importance")

# Left Card
add_content_card(slide8, Inches(0.8), Inches(2.0), Inches(4.0), Inches(4.2), "", [
    "\U0001F50D Interpreting Black Boxes: Extracted feature importance scores to validate classification triggers.",
    "\u2695 Primary Malignancy Drivers: Uniformity of Cell Size and Uniformity of Cell Shape serve as dominant indicators.",
    "\u2696 Secondary Signifiers: Bare Nuclei and Bland Chromatin show moderate relative tracking."
])

# Middle image card
add_visual_card(slide8, "feature_importance.png", Inches(5.0), Inches(2.0), Inches(4.3), Inches(4.2))

# Presenter guide
add_presenter_guide(slide8, "Phase 7: Explainability", "AI Explanations", [
    "\u2192 Tackles 'Black-Box' issues.",
    "\u2192 Verifies clinical indicators physically."
], "By showcasing feature weights, clinicians can confirm that physical tumor indicators align with computer-calculated decisions.")

add_footer(slide8, 8, center_text="FEATURE IMPORTANCES")


# ==========================================
# SLIDE 9: Interactive Diagnostic Web Portal (With Graphics)
# ==========================================
slide9 = prs.slides.add_slide(slide_layout)
set_dark_background(slide9)
add_slide_title(slide9, "Interactive Medical-Grade", "Web Portal Preview")

# Left image card (using Gradio mock UI)
add_visual_card(slide9, "gradio_mockup.png", Inches(0.8), Inches(2.0), Inches(4.3), Inches(4.2))

# Middle Card
add_content_card(slide9, Inches(5.4), Inches(2.0), Inches(4.1), Inches(4.2), "", [
    "\U0001F4BB Engine Framework: Built via Python-based Gradio interfaces ('app.py') for clean, local browser deployment.",
    "\u2699 Frictionless Slide Inputs: Utilizes 9 biopsy attribute sliders ranging from 1 to 10 for easy, rapid trial tests.",
    "\U0001F4C8 Dual Real-time Output: Renders parallel ML estimates side-by-side with Deep Learning risk profiles."
])

# Presenter guide
add_presenter_guide(slide9, "Phase 8: Deployment", "Portal Execution", [
    "\u2192 Instant Gradio web testing.",
    "\u2192 Physician-friendly slider inputs."
], "The clinician changes input variables in real-time, receiving parallel diagnostic reports in under 50 milliseconds.")

add_footer(slide9, 9, center_text="PORTAL DEPLOYMENT")


# ==========================================
# SLIDE 10: Conclusion & Next Steps
# ==========================================
slide10 = prs.slides.add_slide(slide_layout)
set_dark_background(slide10)
add_slide_title(slide10, "Summary &", "Future Integration Vision")

# Draw Timeline line
t_line = slide10.shapes.add_shape(1, Inches(0.8), Inches(4.0), Inches(8.8), Inches(0.04))
t_line.fill.solid()
t_line.fill.fore_color.rgb = TEAL
t_line.line.fill.background()

# Helper for timeline card
def add_timeline_card(slide, left, top, width, height, phase, desc, is_bottom=True):
    # Card Background
    card = slide10.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_DARK
    card.line.color.rgb = TEAL
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    
    # Phase text
    p_phase = tf.paragraphs[0]
    p_phase.alignment = PP_ALIGN.CENTER
    p_phase.text = phase
    p_phase.font.bold = True
    p_phase.font.size = Pt(13)
    p_phase.font.color.rgb = WHITE
    p_phase.font.name = "Arial"
    p_phase.space_after = Pt(4)
    
    # Description text
    p_desc = tf.add_paragraph()
    p_desc.alignment = PP_ALIGN.CENTER
    p_desc.text = desc
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = GREY_TEXT
    p_desc.font.name = "Arial"
    
    # Node coordinate
    node_top = Inches(3.9)
    node_left = left + width / 2 - Inches(0.125)
    
    # Draw circle node on line
    node = slide10.shapes.add_shape(9, node_left, node_top, Inches(0.25), Inches(0.25))
    node.fill.solid()
    node.fill.fore_color.rgb = TEAL
    node.line.color.rgb = BG_DARK
    node.line.width = Pt(2.5)

# Add Timeline Cards (alternating top and bottom)
add_timeline_card(slide10, Inches(0.8), Inches(4.3), Inches(2.1), Inches(2.2), "Phase 1", "Baseline Complete: Achieved accuracy limits exceeding 95.5% on baseline testing.", is_bottom=True)
add_timeline_card(slide10, Inches(3.0), Inches(1.5), Inches(2.1), Inches(2.2), "Phase 2", "Serialization: Package predictive nodes into safe pickle wrappers.", is_bottom=False)
add_timeline_card(slide10, Inches(5.2), Inches(4.3), Inches(2.1), Inches(2.2), "Phase 3", "Clinical Scaling: Launch live validation trials across partner center hospital databases.", is_bottom=True)
add_timeline_card(slide10, Inches(7.4), Inches(1.5), Inches(2.1), Inches(2.2), "Phase 4", "EHR Pipelines: Deploy model endpoints via secure APIs into oncology records.", is_bottom=False)

# Presenter guide
add_presenter_guide(slide10, "Phase 9: Scale", "Hospital Rollout", [
    "\u2192 Future-proof EHR integrations.",
    "\u2192 REST API scalability paths."
], "Our final step is wrapping these validated classifiers into REST API endpoints connecting directly with clinical databases.")

add_footer(slide10, 10, center_text="ROADMAP SCALE")


# ==========================================
# SLIDE 11: Image Sources
# ==========================================
slide11 = prs.slides.add_slide(slide_layout)
set_dark_background(slide11)

# Slide Title
tx_box11 = slide11.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
tf11 = tx_box11.text_frame
p11 = tf11.paragraphs[0]
p11.text = "Image Sources"
p11.font.bold = True
p11.font.size = Pt(28)
p11.font.color.rgb = WHITE
p11.font.name = "Arial"

# Add microscope image thumbnail
slide11.shapes.add_picture("fna_scan.jpg", Inches(0.8), Inches(2.5), Inches(2.0), Inches(1.5))

# Text Box for source URLs
url_box = slide11.shapes.add_textbox(Inches(3.2), Inches(2.5), Inches(9.0), Inches(1.5))
url_tf = url_box.text_frame
url_tf.word_wrap = True

p_url = url_tf.paragraphs[0]
p_url.text = "https://www.scientificarchives.com/public/assets/images/uploads/image-1707331339-2.jpg"
p_url.font.size = Pt(14)
p_url.font.color.rgb = TEAL
p_url.font.name = "Arial"
p_url.space_after = Pt(10)

p_src = url_tf.add_paragraph()
p_src.text = "Source: www.scientificarchives.com"
p_src.font.size = Pt(14)
p_src.font.color.rgb = WHITE
p_src.font.name = "Arial"

add_footer(slide11, 11, total=11, center_text="RESOURCES")

# Save updated presentation
prs.save("breast_cancer_diagnostic_presentation_v3.pptx")
print("Presentation updated successfully with the premium layout and real graphical data!")

